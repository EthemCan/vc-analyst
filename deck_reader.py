"""Utilities for reading and normalizing pitch deck text from .pptx/.pdf files."""

from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path
from typing import List

from pptx import Presentation
from pypdf import PdfReader


@dataclass
class DeckReader:
    """Extracts text from .pptx or .pdf pitch decks into a structured string."""

    min_line_length: int = 2

    def extract_text(self, file_path: str) -> str:
        """
        Extract and clean text from a supported pitch deck format.

        Args:
            file_path: Path to the .pptx or .pdf pitch deck.

        Returns:
            A structured string grouped by slide.

        Raises:
            FileNotFoundError: If the deck path does not exist.
            ValueError: If file extension is invalid or no text can be extracted.
            RuntimeError: If parser fails to read the file.
        """
        path = Path(file_path).expanduser().resolve()
        if not path.exists():
            raise FileNotFoundError(f"Pitch deck not found: {path}")
        suffix = path.suffix.lower()
        if suffix == ".pptx":
            return self._extract_from_pptx(path)
        if suffix == ".pdf":
            return self._extract_from_pdf(path)
        raise ValueError(f"Unsupported file type '{suffix}'. Use .pptx or .pdf.")

    def _extract_from_pptx(self, path: Path) -> str:
        """Extract text from PowerPoint slides."""
        try:
            presentation = Presentation(str(path))
        except Exception as exc:  # pragma: no cover - dependent on corrupt file state
            raise RuntimeError(f"Failed to open PowerPoint deck '{path}': {exc}") from exc

        slide_sections: List[str] = []
        for index, slide in enumerate(presentation.slides, start=1):
            raw_lines = []
            for shape in slide.shapes:
                text = self._shape_text(shape)
                if text:
                    raw_lines.extend(text.splitlines())

            cleaned_lines = self._clean_lines(raw_lines)
            if not cleaned_lines:
                continue

            section = [f"Slide {index}"] + [f"- {line}" for line in cleaned_lines]
            slide_sections.append("\n".join(section))

        if not slide_sections:
            raise ValueError("No readable text found in the provided .pptx file.")
        return "\n\n".join(slide_sections)

    def _extract_from_pdf(self, path: Path) -> str:
        """Extract text from PDF pages."""
        try:
            reader = PdfReader(str(path))
        except Exception as exc:  # pragma: no cover - dependent on corrupt file state
            raise RuntimeError(f"Failed to open PDF deck '{path}': {exc}") from exc

        page_sections: List[str] = []
        for index, page in enumerate(reader.pages, start=1):
            try:
                raw_text = page.extract_text() or ""
            except Exception as exc:  # pragma: no cover - unreadable page edge cases
                raise RuntimeError(
                    f"Failed to read text from PDF page {index} in '{path}': {exc}"
                ) from exc

            cleaned_lines = self._clean_lines(raw_text.splitlines())
            if not cleaned_lines:
                continue

            section = [f"Page {index}"] + [f"- {line}" for line in cleaned_lines]
            page_sections.append("\n".join(section))

        if not page_sections:
            raise ValueError("No readable text found in the provided .pdf file.")
        return "\n\n".join(page_sections)

    def _shape_text(self, shape: object) -> str:
        """Safely fetch shape text if available."""
        if hasattr(shape, "has_text_frame") and shape.has_text_frame:
            return (shape.text or "").strip()
        if hasattr(shape, "text"):
            return (shape.text or "").strip()
        return ""

    def _clean_lines(self, lines: List[str]) -> List[str]:
        """Normalize whitespace and filter out noisy lines."""
        normalized: List[str] = []
        for line in lines:
            cleaned = " ".join(line.split())
            if len(cleaned) >= self.min_line_length:
                normalized.append(cleaned)
        return normalized
